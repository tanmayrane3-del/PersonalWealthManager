"""
Investor Pitch Deck Generator for Personal Wealth Manager
Run: python generate_pitch_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── Color Palette (Wealth / Fintech theme) ────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # Deep navy background
GOLD        = RGBColor(0xF5, 0xA6, 0x23)   # Gold accent
LIGHT_GOLD  = RGBColor(0xFF, 0xD8, 0x80)   # Soft gold
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF0, 0xF4, 0xFF)
GRAY        = RGBColor(0xB0, 0xB8, 0xCC)
DARK_GRAY   = RGBColor(0x33, 0x3C, 0x50)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
TEAL        = RGBColor(0x1A, 0xBC, 0x9C)
CARD_BG     = RGBColor(0x16, 0x25, 0x4D)   # Slightly lighter navy for cards

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def fill_bg(slide, color):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width > Pt(0) else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Calibri", word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=16, color=WHITE, align=PP_ALIGN.LEFT,
                       font_name="Calibri", line_spacing=None, bold_first=False):
    """Each item in lines is either a str or (str, dict) for per-line overrides."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, opts = line
        else:
            text, opts = line, {}
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = opts.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", font_size))
        run.font.bold = opts.get("bold", bold_first and i == 0)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", color)
        run.font.name = opts.get("font", font_name)
    return txBox


def divider(slide, y, color=GOLD, width_pct=0.9):
    """Horizontal accent line."""
    w = SLIDE_W * width_pct
    x = (SLIDE_W - w) / 2
    add_rect(slide, x, y, w, Pt(2), fill_color=color)


def slide_number(slide, num, total):
    add_text(slide, f"{num} / {total}",
             SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4),
             Inches(1.0), Inches(0.3),
             font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def section_tag(slide, label):
    """Small gold pill in top-left."""
    add_rect(slide, Inches(0.4), Inches(0.25), Inches(2.4), Inches(0.32),
             fill_color=GOLD)
    add_text(slide, label.upper(),
             Inches(0.4), Inches(0.25), Inches(2.4), Inches(0.32),
             font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def footer_bar(slide):
    """Thin gold bar at bottom."""
    add_rect(slide, 0, SLIDE_H - Inches(0.08),
             SLIDE_W, Inches(0.08), fill_color=GOLD)


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Full-height left panel (gold accent strip)
    add_rect(s, 0, 0, Inches(0.15), SLIDE_H, fill_color=GOLD)

    # Right accent strip
    add_rect(s, SLIDE_W - Inches(0.15), 0, Inches(0.15), SLIDE_H, fill_color=GOLD)

    # Diagonal decorative block (simulated with a wide rotated rectangle — skip for simplicity)
    # Semi-transparent overlay card
    add_rect(s, Inches(0.5), Inches(1.6), Inches(8.5), Inches(4.3), fill_color=CARD_BG)

    # App name
    add_text(s, "WealthPulse",
             Inches(0.9), Inches(1.9), Inches(8.0), Inches(1.1),
             font_size=52, bold=True, color=GOLD, align=PP_ALIGN.LEFT,
             font_name="Calibri Light")

    # Tagline
    add_text(s, "Your Complete Financial Command Centre",
             Inches(0.9), Inches(2.9), Inches(8.0), Inches(0.6),
             font_size=24, color=WHITE, align=PP_ALIGN.LEFT)

    # Sub-tagline
    add_text(s, "Track every rupee. Grow every asset. Own your future.",
             Inches(0.9), Inches(3.55), Inches(8.0), Inches(0.5),
             font_size=16, italic=True, color=LIGHT_GOLD, align=PP_ALIGN.LEFT)

    divider(s, Inches(4.2), width_pct=0.62)

    # Pitch meta
    add_multiline_text(s, [
        ("Seed Round Investment Opportunity", {"size": 13, "bold": True, "color": WHITE}),
        ("Personal Finance & Wealth Management  |  India  |  2026", {"size": 12, "color": GRAY}),
    ], Inches(0.9), Inches(4.35), Inches(8.0), Inches(0.8))

    # Right side — key stats bubble (decorative)
    add_rect(s, Inches(10.0), Inches(1.5), Inches(2.9), Inches(4.5), fill_color=CARD_BG)
    add_text(s, "SNAPSHOT", Inches(10.0), Inches(1.7), Inches(2.9), Inches(0.4),
             font_size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    stats = [
        ("6+", "Asset Classes"),
        ("21", "API Modules"),
        ("AI", "SMS Parsing"),
        ("Live", "Market Data"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = Inches(2.2) + i * Inches(0.85)
        add_text(s, val, Inches(10.1), y, Inches(1.2), Inches(0.45),
                 font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, lbl, Inches(11.2), y + Inches(0.1), Inches(1.6), Inches(0.35),
                 font_size=11, color=GRAY, align=PP_ALIGN.LEFT)

    footer_bar(s)
    return s


def slide_02_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "The Problem")
    footer_bar(s)

    add_text(s, "India's Investors Are Flying Blind",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.7),
             font_size=34, bold=True, color=WHITE)
    divider(s, Inches(1.55), width_pct=0.94)

    problems = [
        ("💸  Fragmented Wealth",
         "Stocks in Zerodha, mutual funds in CAMS, gold in a locker, EMIs in a spreadsheet. No single view of true net worth."),
        ("📉  No Growth Visibility",
         "Investors can't see which assets are growing, which are dragging — making rebalancing guesswork."),
        ("🧾  Manual Expense Chaos",
         "Tracking daily spends requires manual data entry. SMS alerts are ignored, money leaks unnoticed."),
        ("📊  Zero Projections",
         "No tool tells an Indian retail investor: \"At your current pace, your net worth will be ₹X in 5 years.\""),
    ]

    cols = [(Inches(0.4), Inches(1.75)), (Inches(6.8), Inches(1.75)),
            (Inches(0.4), Inches(4.35)), (Inches(6.8), Inches(4.35))]

    for i, (title, body) in enumerate(problems):
        x, y = cols[i]
        add_rect(s, x, y, Inches(6.1), Inches(2.3), fill_color=CARD_BG)
        add_rect(s, x, y, Inches(0.1), Inches(2.3), fill_color=GOLD)
        add_text(s, title, x + Inches(0.25), y + Inches(0.15), Inches(5.7), Inches(0.45),
                 font_size=15, bold=True, color=GOLD)
        add_text(s, body, x + Inches(0.25), y + Inches(0.65), Inches(5.7), Inches(1.5),
                 font_size=13, color=WHITE)

    slide_number(s, 2, 12)


def slide_03_solution(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "The Solution")
    footer_bar(s)

    add_text(s, "WealthPulse — One App for Your Entire Financial Life",
             Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.7),
             font_size=30, bold=True, color=WHITE)
    divider(s, Inches(1.55), width_pct=0.94)

    add_text(s, "A native Android wealth management platform that automatically aggregates, tracks, and projects every financial asset and liability — powered by live market data and AI.",
             Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.65),
             font_size=15, italic=True, color=LIGHT_GOLD)

    pillars = [
        ("🏦", "Aggregate", "Connect stocks, mutual funds, gold, real estate, and loans in one dashboard"),
        ("📈", "Track", "Live prices, day P&L, CAGR per asset class — updated during market hours"),
        ("🤖", "Automate", "AI-powered SMS parsing auto-captures every bank transaction, zero manual entry"),
        ("🔮", "Project", "5-year net worth projections with CAGR-based growth forecasts per asset class"),
    ]

    for i, (icon, title, desc) in enumerate(pillars):
        x = Inches(0.35) + i * Inches(3.2)
        add_rect(s, x, Inches(2.5), Inches(3.0), Inches(3.9), fill_color=CARD_BG)
        add_rect(s, x, Inches(2.5), Inches(3.0), Inches(0.12), fill_color=GOLD)
        add_text(s, icon, x + Inches(0.1), Inches(2.72), Inches(2.8), Inches(0.55),
                 font_size=28, align=PP_ALIGN.CENTER)
        add_text(s, title, x + Inches(0.1), Inches(3.35), Inches(2.8), Inches(0.45),
                 font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, desc, x + Inches(0.15), Inches(3.85), Inches(2.75), Inches(1.4),
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number(s, 3, 12)


def slide_04_product(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Product")
    footer_bar(s)

    add_text(s, "Feature-Complete MVP",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.6),
             font_size=34, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    features = [
        ("Net Worth Dashboard", "Real-time aggregation of all assets & liabilities with 5Y projection"),
        ("Stocks (Zerodha OAuth)", "Live holdings sync via Kite Connect API — P&L, CAGR 1Y/3Y/5Y"),
        ("Mutual Funds (CAS Import)", "PDF parsing → AMFI NAV sync → XIRR return calculation"),
        ("Gold & Silver", "Live MCX rates, quantity tracking by purity (24K/22K)"),
        ("Physical Assets", "Real estate & vehicles with WDV depreciation calculations"),
        ("Liabilities & Loans", "EMI calculation, principal balance, payoff projections"),
        ("AI Expense Auto-Capture", "Bank SMS → Claude AI → categorised transaction, automatic"),
        ("Transaction Management", "Income/expense CRUD with custom categories, sources, recipients"),
        ("Growth Projections", "Per-class CAGR model → full 1Y / 3Y / 5Y net worth forecast"),
        ("Cron Automation", "Market sync, daily snapshots, weekly CAGR recalc — all scheduled"),
    ]

    mid = len(features) // 2
    for col, chunk in enumerate([features[:mid], features[mid:]]):
        x = Inches(0.4) + col * Inches(6.55)
        for row, (title, desc) in enumerate(chunk):
            y = Inches(1.65) + row * Inches(1.05)
            add_rect(s, x, y, Inches(6.2), Inches(0.93), fill_color=CARD_BG)
            add_rect(s, x, y, Inches(0.07), Inches(0.93), fill_color=GREEN)
            add_text(s, "✓  " + title, x + Inches(0.2), y + Inches(0.06), Inches(5.8), Inches(0.38),
                     font_size=13, bold=True, color=GREEN)
            add_text(s, desc, x + Inches(0.2), y + Inches(0.48), Inches(5.8), Inches(0.38),
                     font_size=11, color=GRAY)

    slide_number(s, 4, 12)


def slide_05_tech(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Technology")
    footer_bar(s)

    add_text(s, "Modern, Scalable Tech Stack",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.6),
             font_size=34, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    # Architecture flow boxes
    layers = [
        ("Android App (Kotlin)", "MVVM · Hilt DI · Retrofit\nRoom DB · WorkManager\nMPAndroidChart · Coroutines", GOLD),
        ("REST API (Node.js)", "Express.js 5 · 21 Modules\nJWT Sessions · Middleware\nPDF & SMS Parsers", RGBColor(0x1A, 0xBC, 0x9C)),
        ("Database (Supabase)", "PostgreSQL · PgBouncer\nSnapshots · CAGR Cache\nMetal Rates Cache", RGBColor(0x9B, 0x59, 0xB6)),
        ("External APIs", "Zerodha Kite Connect\nAMFI · Yahoo Finance\nClaude AI · MCX Rates", RGBColor(0xE7, 0x4C, 0x3C)),
    ]

    for i, (title, body, color) in enumerate(layers):
        x = Inches(0.35) + i * Inches(3.2)
        add_rect(s, x, Inches(1.7), Inches(3.0), Inches(4.5), fill_color=CARD_BG)
        add_rect(s, x, Inches(1.7), Inches(3.0), Inches(0.15), fill_color=color)
        add_text(s, title, x + Inches(0.12), Inches(2.0), Inches(2.8), Inches(0.45),
                 font_size=14, bold=True, color=color)
        add_text(s, body, x + Inches(0.12), Inches(2.55), Inches(2.8), Inches(3.3),
                 font_size=12, color=WHITE)

    # Arrow connectors (simple text arrows)
    for i in range(3):
        x = Inches(3.25) + i * Inches(3.2)
        add_text(s, "→", x, Inches(3.6), Inches(0.25), Inches(0.5),
                 font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Deployment note
    add_rect(s, Inches(0.35), Inches(6.45), Inches(12.6), Inches(0.65), fill_color=DARK_GRAY)
    add_text(s, "🚀  Deployed:  Backend on Render.com  |  Database on Supabase  |  Android App (dev-ready for Play Store)",
             Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.55),
             font_size=12, color=WHITE)

    slide_number(s, 5, 12)


def slide_06_market(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Market Opportunity")
    footer_bar(s)

    add_text(s, "A $1.5T Market — Underserved by Digital Wealth Tools",
             Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.65),
             font_size=30, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    # Big numbers
    stats = [
        ("50M+", "Retail investors\nin India (2026)", GOLD),
        ("₹50 Lakh Cr", "Equity AUM in\nIndian markets", TEAL),
        ("8× Growth", "Demat accounts\nin 5 years", GREEN),
        ("$6B+", "India Fintech\nfunding in 2023", RGBColor(0xE7, 0x4C, 0x3C)),
    ]

    for i, (val, lbl, color) in enumerate(stats):
        x = Inches(0.35) + i * Inches(3.2)
        add_rect(s, x, Inches(1.65), Inches(3.0), Inches(2.1), fill_color=CARD_BG)
        add_text(s, val, x, Inches(1.8), Inches(3.0), Inches(0.85),
                 font_size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, lbl, x, Inches(2.65), Inches(3.0), Inches(0.9),
                 font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(s, "Why Now",
             Inches(0.5), Inches(4.0), Inches(12.0), Inches(0.45),
             font_size=18, bold=True, color=GOLD)

    why_now = [
        "📱  India has 700M+ smartphones; mobile-first finance is the default for Gen Z & Millennials.",
        "📈  Post-COVID DIY investing surge — millions managing their own portfolios for the first time.",
        "🏦  Zerodha alone has 14M active users; no integrated wealth tracker serves them natively.",
        "🤖  AI (LLMs) now makes automated transaction parsing viable at near-zero marginal cost.",
        "📋  SEBI's push for financial literacy creates policy tailwind for wealth management apps.",
    ]
    add_multiline_text(s, why_now,
                       Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.7),
                       font_size=13, color=WHITE)

    slide_number(s, 6, 12)


def slide_07_traction(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Traction & Status")
    footer_bar(s)

    add_text(s, "MVP Built. Backend Live. Ready to Scale.",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.6),
             font_size=34, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    built = [
        "✅  Full-stack MVP with 10 screens and 21 backend API modules",
        "✅  Zerodha (India's #1 broker) OAuth integration live and tested",
        "✅  AI-powered SMS expense capture (Claude API) working end-to-end",
        "✅  CAS PDF import for mutual funds fully operational",
        "✅  Real-time gold/silver rates, stock prices, AMFI NAV sync",
        "✅  5-year net worth projection engine with per-asset CAGR model",
        "✅  Automated cron jobs: daily snapshots, market sync, weekly CAGR",
        "✅  Backend deployed on Render · Database live on Supabase PostgreSQL",
    ]

    add_multiline_text(s, built,
                       Inches(0.4), Inches(1.65), Inches(7.8), Inches(4.8),
                       font_size=13, color=WHITE)

    # Milestone timeline on the right
    add_rect(s, Inches(8.6), Inches(1.55), Inches(4.4), Inches(5.5), fill_color=CARD_BG)
    add_text(s, "MILESTONES", Inches(8.6), Inches(1.7), Inches(4.4), Inches(0.4),
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    milestones = [
        ("Q3 2025", "Architecture design & DB schema", GRAY),
        ("Q4 2025", "Core modules: Auth, Transactions, Stocks", GRAY),
        ("Q1 2026", "Zerodha OAuth, Mutual Funds, Metals", GOLD),
        ("Mar 2026", "Net Worth projections + AI SMS ✓ NOW", GREEN),
        ("Q2 2026", "Play Store launch + beta users", LIGHT_GOLD),
        ("Q3 2026", "Premium tier + monetisation", LIGHT_GOLD),
    ]
    for i, (date, event, color) in enumerate(milestones):
        y = Inches(2.25) + i * Inches(0.72)
        add_rect(s, Inches(8.75), y, Inches(0.08), Inches(0.5), fill_color=color)
        add_text(s, date, Inches(8.9), y + Inches(0.02), Inches(1.3), Inches(0.35),
                 font_size=10, bold=True, color=color)
        add_text(s, event, Inches(10.25), y + Inches(0.02), Inches(2.65), Inches(0.45),
                 font_size=11, color=WHITE)

    slide_number(s, 7, 12)


def slide_08_business_model(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Business Model")
    footer_bar(s)

    add_text(s, "Multiple Revenue Streams on a SaaS Foundation",
             Inches(0.5), Inches(0.7), Inches(12.5), Inches(0.6),
             font_size=30, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    streams = [
        ("💎  Premium Subscriptions",
         "₹199–499/month",
         "Advanced analytics, tax reports (LTCG), portfolio rebalancing alerts, export features",
         GOLD),
        ("🤝  Brokerage Partnerships",
         "Revenue Share",
         "Affiliate commissions from Zerodha, Groww, Kuvera — referral on new account opens",
         TEAL),
        ("🏢  Enterprise / Family Plans",
         "₹999+/month",
         "Multi-user household wealth tracking, shared dashboards, wealth advisor portal",
         GREEN),
        ("📊  Advisory Marketplace",
         "Commission-based",
         "In-app SEBI RIA consultation booking, premium investment reports, goal planning",
         RGBColor(0x9B, 0x59, 0xB6)),
        ("📱  White-label SDK",
         "B2B Licensing",
         "Net worth & portfolio engine licensed to banks, NBFCs, and insurance platforms",
         RGBColor(0xE7, 0x4C, 0x3C)),
        ("🔑  API Access",
         "Usage-based",
         "Wealth data API for fintech developers — CAGR engine, SMS parser, CAS importer",
         RGBColor(0xF3, 0x9C, 0x12)),
    ]

    for i, (title, revenue, desc, color) in enumerate(streams):
        row, col = divmod(i, 3)
        x = Inches(0.35) + col * Inches(4.3)
        y = Inches(1.65) + row * Inches(2.6)
        add_rect(s, x, y, Inches(4.1), Inches(2.4), fill_color=CARD_BG)
        add_rect(s, x, y, Inches(4.1), Inches(0.1), fill_color=color)
        add_text(s, title, x + Inches(0.15), y + Inches(0.22), Inches(3.8), Inches(0.45),
                 font_size=13, bold=True, color=color)
        add_text(s, revenue, x + Inches(0.15), y + Inches(0.7), Inches(3.8), Inches(0.35),
                 font_size=12, bold=True, color=WHITE)
        add_text(s, desc, x + Inches(0.15), y + Inches(1.1), Inches(3.8), Inches(1.15),
                 font_size=11, color=GRAY)

    slide_number(s, 8, 12)


def slide_09_competition(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Competitive Landscape")
    footer_bar(s)

    add_text(s, "We Do What No One Else Does — End-to-End",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.6),
             font_size=34, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    headers = ["Feature", "WealthPulse", "Zerodha Kite", "ET Money", "Paytm Money", "Groww"]
    rows = [
        ["Stocks (Live Holdings)",        "✅", "✅", "❌", "❌", "✅"],
        ["Mutual Funds (CAS Import)",     "✅", "❌", "✅", "✅", "✅"],
        ["Gold & Silver Tracking",        "✅", "❌", "❌", "❌", "❌"],
        ["Physical Assets / Real Estate", "✅", "❌", "❌", "❌", "❌"],
        ["Liabilities / Loan Tracking",   "✅", "❌", "❌", "❌", "❌"],
        ["Net Worth Projection (5Y)",     "✅", "❌", "❌", "❌", "❌"],
        ["AI-Powered SMS Expense Capture","✅", "❌", "❌", "⚠️", "❌"],
        ["Unified Net Worth Dashboard",   "✅", "❌", "⚠️", "❌", "❌"],
    ]

    col_widths = [Inches(3.3), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65)]
    col_x = [Inches(0.3)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    row_h = Inches(0.6)
    header_y = Inches(1.6)

    # Header row
    for j, (h, w, x) in enumerate(zip(headers, col_widths, col_x)):
        bg = GOLD if j == 1 else DARK_GRAY
        fc = NAVY if j == 1 else GRAY
        add_rect(s, x, header_y, w, Inches(0.5), fill_color=bg)
        add_text(s, h, x + Inches(0.05), header_y + Inches(0.06), w - Inches(0.1), Inches(0.38),
                 font_size=11, bold=True, color=fc, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        y = header_y + Inches(0.55) + i * row_h
        bg = CARD_BG if i % 2 == 0 else NAVY
        for j, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
            cell_bg = RGBColor(0x1A, 0x2F, 0x5A) if j == 1 else bg
            add_rect(s, x, y, w, row_h - Inches(0.05), fill_color=cell_bg)
            color = GREEN if cell == "✅" else (RGBColor(0xE7, 0x4C, 0x3C) if cell == "❌" else GOLD)
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            add_text(s, cell, x + Inches(0.05), y + Inches(0.1), w - Inches(0.08), Inches(0.4),
                     font_size=12, color=color, align=align,
                     bold=(j == 1))

    slide_number(s, 9, 12)


def slide_10_roadmap(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "Roadmap")
    footer_bar(s)

    add_text(s, "From MVP to Market Leader — 18-Month Plan",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.6),
             font_size=34, bold=True, color=WHITE)
    divider(s, Inches(1.45), width_pct=0.94)

    phases = [
        ("Phase 1\nQ2 2026", [
            "Play Store launch (public beta)",
            "Push notifications & price alerts",
            "Offline mode with sync",
            "Onboarding flow & tutorial",
            "500 beta users target",
        ], GOLD),
        ("Phase 2\nQ3 2026", [
            "Premium subscription tier",
            "LTCG tax report export",
            "Portfolio rebalancing suggestions",
            "Crypto holdings support",
            "Zerodha affiliate integration live",
        ], TEAL),
        ("Phase 3\nQ4 2026", [
            "Family / multi-user accounts",
            "SEBI RIA advisory marketplace",
            "Web dashboard (React)",
            "Insurance tracking module",
            "B2B white-label pilot",
        ], GREEN),
        ("Phase 4\nQ1-2 2027", [
            "10,000 paid subscribers",
            "API access tier launch",
            "iOS app development",
            "Series A fundraise",
            "Pan-India marketing push",
        ], RGBColor(0x9B, 0x59, 0xB6)),
    ]

    for i, (phase, items, color) in enumerate(phases):
        x = Inches(0.35) + i * Inches(3.2)
        add_rect(s, x, Inches(1.65), Inches(3.0), Inches(5.4), fill_color=CARD_BG)
        add_rect(s, x, Inches(1.65), Inches(3.0), Inches(0.12), fill_color=color)
        add_text(s, phase, x + Inches(0.1), Inches(1.85), Inches(2.85), Inches(0.7),
                 font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            y = Inches(2.65) + j * Inches(0.8)
            add_text(s, "▸  " + item, x + Inches(0.15), y, Inches(2.8), Inches(0.72),
                     font_size=12, color=WHITE)

    slide_number(s, 10, 12)


def slide_11_ask(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    section_tag(s, "The Ask")
    footer_bar(s)

    add_text(s, "Seed Round: ₹1.5 Crore",
             Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.7),
             font_size=38, bold=True, color=GOLD)
    divider(s, Inches(1.5), width_pct=0.94)

    # Use of funds
    add_text(s, "Use of Funds", Inches(0.5), Inches(1.65), Inches(6.0), Inches(0.45),
             font_size=18, bold=True, color=WHITE)

    funds = [
        ("40%  ₹60L", "Product & Engineering", "2 senior developers, iOS port, web dashboard, infra scaling", GOLD),
        ("25%  ₹37.5L", "Go-To-Market", "Play Store launch, content marketing, influencer partnerships", TEAL),
        ("20%  ₹30L", "Compliance & Legal", "SEBI RIA registration, data privacy, security audits", GREEN),
        ("15%  ₹22.5L", "Operations & Runway", "18-month founder + team runway, cloud infra, API costs", GRAY),
    ]

    for i, (pct, title, desc, color) in enumerate(funds):
        y = Inches(2.2) + i * Inches(1.05)
        add_rect(s, Inches(0.4), y, Inches(6.3), Inches(0.95), fill_color=CARD_BG)
        add_rect(s, Inches(0.4), y, Inches(0.1), Inches(0.95), fill_color=color)
        add_text(s, pct, Inches(0.65), y + Inches(0.08), Inches(1.5), Inches(0.4),
                 font_size=13, bold=True, color=color)
        add_text(s, title, Inches(2.3), y + Inches(0.08), Inches(4.2), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)
        add_text(s, desc, Inches(0.65), y + Inches(0.53), Inches(5.9), Inches(0.35),
                 font_size=11, color=GRAY)

    # What we offer
    add_rect(s, Inches(7.1), Inches(1.65), Inches(5.8), Inches(5.0), fill_color=CARD_BG)
    add_text(s, "What Investors Get", Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.45),
             font_size=16, bold=True, color=GOLD)
    offer = [
        "🎯  Equity stake in India's most complete personal wealth app",
        "📊  Working MVP — zero technical risk on core product",
        "🔑  Zerodha integration gives immediate access to 14M users",
        "📈  3 proven revenue streams, multiple expansion paths",
        "🇮🇳  Riding India's retail investor boom (50M+ and growing)",
        "🤖  AI + automation moat — hard to replicate quickly",
        "📅  18-month runway to Series A milestones",
    ]
    add_multiline_text(s, offer,
                       Inches(7.2), Inches(2.4), Inches(5.5), Inches(4.0),
                       font_size=12.5, color=WHITE)

    slide_number(s, 11, 12)


def slide_12_closing(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    add_rect(s, 0, 0, Inches(0.15), SLIDE_H, fill_color=GOLD)
    add_rect(s, SLIDE_W - Inches(0.15), 0, Inches(0.15), SLIDE_H, fill_color=GOLD)

    add_text(s, "WealthPulse",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.1),
             font_size=56, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
             font_name="Calibri Light")

    add_text(s, "Every Indian Investor Deserves to See Their True Wealth",
             Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.65),
             font_size=22, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    divider(s, Inches(3.3), width_pct=0.5)

    add_text(s, "Let's build India's #1 wealth management platform together.",
             Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.55),
             font_size=18, color=LIGHT_GOLD, align=PP_ALIGN.CENTER)

    # Contact placeholder
    add_rect(s, Inches(3.8), Inches(4.3), Inches(5.7), Inches(2.2), fill_color=CARD_BG)
    add_text(s, "Get in Touch",
             Inches(3.8), Inches(4.5), Inches(5.7), Inches(0.45),
             font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_multiline_text(s, [
        ("📧  [your.email@domain.com]", {"align": PP_ALIGN.CENTER, "color": WHITE, "size": 13}),
        ("📱  [+91-XXXXX-XXXXX]", {"align": PP_ALIGN.CENTER, "color": WHITE, "size": 13}),
        ("🌐  [github / portfolio link]", {"align": PP_ALIGN.CENTER, "color": GRAY, "size": 12}),
    ], Inches(3.8), Inches(5.05), Inches(5.7), Inches(1.3))

    footer_bar(s)
    add_text(s, "12 / 12", SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4),
             Inches(1.0), Inches(0.3),
             font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs)
    print("  1/12  Title")
    slide_02_problem(prs)
    print("  2/12  Problem")
    slide_03_solution(prs)
    print("  3/12  Solution")
    slide_04_product(prs)
    print("  4/12  Product Features")
    slide_05_tech(prs)
    print("  5/12  Tech Stack")
    slide_06_market(prs)
    print("  6/12  Market Opportunity")
    slide_07_traction(prs)
    print("  7/12  Traction & Status")
    slide_08_business_model(prs)
    print("  8/12  Business Model")
    slide_09_competition(prs)
    print("  9/12  Competition")
    slide_10_roadmap(prs)
    print(" 10/12  Roadmap")
    slide_11_ask(prs)
    print(" 11/12  The Ask")
    slide_12_closing(prs)
    print(" 12/12  Closing")

    out_path = os.path.join(os.path.dirname(__file__), "WealthPulse_InvestorPitch.pptx")
    prs.save(out_path)
    print(f"\n✅  Saved: {out_path}")


if __name__ == "__main__":
    main()
